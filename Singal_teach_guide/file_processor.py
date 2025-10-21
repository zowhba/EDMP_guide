import os
import PyPDF2
import docx
import pptx
import openpyxl
import logging
from typing import Dict, Any, Optional
import mimetypes

logger = logging.getLogger(__name__)


class FileProcessor:
    def __init__(self):
        self.supported_formats = {
            'pdf': self._process_pdf,
            'docx': self._process_docx,
            'doc': self._process_docx,
            'pptx': self._process_pptx,
            'ppt': self._process_pptx,
            'xlsx': self._process_xlsx,
            'xls': self._process_xlsx,
            'txt': self._process_text,
            'md': self._process_text,
            'csv': self._process_text
        }
    
    def process_file(self, file_path: str, original_filename: str = None) -> Dict[str, Any]:
        """파일을 처리하여 텍스트 추출"""
        try:
            if not os.path.exists(file_path):
                raise FileNotFoundError(f"파일을 찾을 수 없습니다: {file_path}")
            
            # 파일 확장자 확인
            filename = original_filename or os.path.basename(file_path)
            ext = filename.split('.')[-1].lower()
            
            if ext not in self.supported_formats:
                return {
                    'source': filename,
                    'status': 'error',
                    'error': f"지원하지 않는 파일 형식: {ext}"
                }
            
            # 파일 처리
            processor = self.supported_formats[ext]
            content = processor(file_path)
            
            return {
                'source': filename,
                'title': filename,
                'content': content,
                'type': f'file_{ext}',
                'status': 'success'
            }
            
        except Exception as e:
            logger.error(f"파일 처리 실패 {filename}: {str(e)}")
            return {
                'source': filename or file_path,
                'status': 'error',
                'error': f"파일 처리 실패: {str(e)}"
            }
    
    def _process_pdf(self, file_path: str) -> str:
        """PDF 파일 처리"""
        text_content = []
        
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text = page.extract_text()
                if text.strip():
                    text_content.append(f"--- 페이지 {page_num + 1} ---\n{text}")
        
        return '\n\n'.join(text_content)
    
    def _process_docx(self, file_path: str) -> str:
        """Word 문서 처리"""
        doc = docx.Document(file_path)
        text_content = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text)
        
        # 표 내용도 추출
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text_content.append(' | '.join(row_text))
        
        return '\n\n'.join(text_content)
    
    def _process_pptx(self, file_path: str) -> str:
        """PowerPoint 파일 처리"""
        prs = pptx.Presentation(file_path)
        text_content = []
        
        for slide_num, slide in enumerate(prs.slides):
            slide_text = [f"--- 슬라이드 {slide_num + 1} ---"]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if len(slide_text) > 1:  # 슬라이드 번호 외에 내용이 있는 경우
                text_content.append('\n'.join(slide_text))
        
        return '\n\n'.join(text_content)
    
    def _process_xlsx(self, file_path: str) -> str:
        """Excel 파일 처리"""
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        text_content = []
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_text = [f"--- 시트: {sheet_name} ---"]
            
            for row in sheet.iter_rows(values_only=True):
                row_values = [str(cell) if cell is not None else '' for cell in row]
                if any(row_values):  # 비어있지 않은 행만
                    sheet_text.append(' | '.join(row_values))
            
            if len(sheet_text) > 1:  # 시트 이름 외에 내용이 있는 경우
                text_content.append('\n'.join(sheet_text))
        
        wb.close()
        return '\n\n'.join(text_content)
    
    def _process_text(self, file_path: str) -> str:
        """텍스트 파일 처리"""
        encodings = ['utf-8', 'cp949', 'euc-kr', 'latin-1']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as file:
                    return file.read()
            except UnicodeDecodeError:
                continue
        
        # 모든 인코딩이 실패한 경우
        raise ValueError("파일 인코딩을 감지할 수 없습니다.")
    
    def get_supported_formats(self) -> list:
        """지원하는 파일 형식 목록 반환"""
        return list(self.supported_formats.keys())
